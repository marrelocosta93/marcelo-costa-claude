"""
pptx_components.py — Biblioteca de componentes visuais para python-pptx
Nivel consultoria McKinsey/Bain/BCG

Uso:
    import sys, os
    sys.path.insert(0, os.path.expanduser('~/.claude/lib'))
    # ou para D:/.claude/lib:
    sys.path.insert(0, 'D:/.claude/lib')
    from pptx_components import *
"""

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

# Namespaces
NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'


def _build_rPr(ns, lang="pt-BR", bold=False, size_pt=None, font="Arial", color_hex=None):
    """Helper que respeita ordem OOXML: solidFill antes de latin/cs."""
    attribs = {'lang': lang, 'dirty': '0'}
    if bold:
        attribs['b'] = '1'
    if size_pt is not None:
        attribs['sz'] = str(int(size_pt * 100))
    rPr = etree.Element(f'{{{NS_A}}}rPr', attrib=attribs)
    if color_hex:
        sf = etree.SubElement(rPr, f'{{{NS_A}}}solidFill')
        etree.SubElement(sf, f'{{{NS_A}}}srgbClr', attrib={'val': color_hex})
    etree.SubElement(rPr, f'{{{NS_A}}}latin', attrib={'typeface': font})
    etree.SubElement(rPr, f'{{{NS_A}}}cs', attrib={'typeface': font})
    return rPr


def extract_slide_texts(slide):
    """Extrair todos os textos do slide antes de limpar. NUNCA perder conteudo."""
    texts = {}
    for shape in slide.shapes:
        if shape.has_text_frame:
            key = f'ph_{shape.placeholder_format.idx}' if shape.is_placeholder else shape.name
            texts[key] = shape.text_frame.text
    return texts


def clear_slide_shapes(slide):
    """Limpar shapes do slide para reconstrucao composicional."""
    spTree = slide.shapes._spTree
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            sp = shape._element
            spTree.remove(sp)


def add_textbox_styled(slide, left, top, width, height, text,
                       font_pt=14, font_name="Arial", color="333333",
                       bold=False, alignment=PP_ALIGN.LEFT):
    """Textbox com estilo configuravel."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_pt)
    run.font.name = font_name
    run.font.color.rgb = RGBColor.from_string(color)
    run.font.bold = bold
    return txBox


def add_badge(slide, x, y, num, color="E89232", size=Inches(0.6)):
    """Numero em circulo colorido."""
    shape = slide.shapes.add_shape(
        9,  # MSO_SHAPE.OVAL
        x, y, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(color)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(num)
    run.font.size = Pt(18)
    run.font.name = "Arial"
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run.font.bold = True
    return shape


def add_card(slide, x, y, w, h, title="", body="",
             bg_color="F5F5F5", border_color="E0E0E0",
             title_color="333333", body_color="666666",
             title_pt=14, body_pt=12):
    """Bloco de conteudo com background e borda."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(bg_color)
    shape.line.color.rgb = RGBColor.from_string(border_color)
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)

    if title:
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(title_pt)
        run.font.name = "Arial"
        run.font.color.rgb = RGBColor.from_string(title_color)
        run.font.bold = True

    if body:
        p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = body
        run.font.size = Pt(body_pt)
        run.font.name = "Arial"
        run.font.color.rgb = RGBColor.from_string(body_color)

    return shape


def add_banner(slide, x, y, w, h, text, color="E89232", font_color="FFFFFF"):
    """Header colorido de coluna."""
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(color)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(12)
    run.font.name = "Arial"
    run.font.color.rgb = RGBColor.from_string(font_color)
    run.font.bold = True
    return shape


def add_quote_box(slide, x, y, w, h, quote_text, attribution="",
                  accent_color="E89232", bg_color="F9F3EB"):
    """Citacao com accent bar lateral."""
    # Accent bar
    bar = slide.shapes.add_shape(1, x, y, Inches(0.08), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor.from_string(accent_color)
    bar.line.fill.background()

    # Quote body
    body = slide.shapes.add_shape(1, x + Inches(0.12), y, w - Inches(0.12), h)
    body.fill.solid()
    body.fill.fore_color.rgb = RGBColor.from_string(bg_color)
    body.line.fill.background()

    tf = body.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.1)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f'"{quote_text}"'
    run.font.size = Pt(13)
    run.font.name = "Arial"
    run.font.color.rgb = RGBColor.from_string("333333")
    run.font.italic = True

    if attribution:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(6)
        run2 = p2.add_run()
        run2.text = f"— {attribution}"
        run2.font.size = Pt(10)
        run2.font.name = "Arial"
        run2.font.color.rgb = RGBColor.from_string("888888")

    return body


def add_hline(slide, x, y, w, color="E0E0E0", height=Pt(1)):
    """Separador horizontal."""
    shape = slide.shapes.add_shape(1, x, y, w, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(color)
    shape.line.fill.background()
    return shape


def add_vbar(slide, x, y, h, color="E89232", width=Inches(0.06)):
    """Bullet visual vertical."""
    shape = slide.shapes.add_shape(1, x, y, width, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(color)
    shape.line.fill.background()
    return shape


def add_accent_underline(slide, x, y, w, color="E89232"):
    """Decoracao abaixo de titulos."""
    return add_hline(slide, x, y, w, color=color, height=Pt(3))


def add_numbered_item(slide, x, y, num, title, desc, accent_color="E89232"):
    """Item numerado com badge + titulo + descricao."""
    add_badge(slide, x, y, num, color=accent_color)
    add_textbox_styled(slide, x + Inches(0.8), y, Inches(10), Inches(0.4),
                       title, font_pt=16, bold=True)
    if desc:
        add_textbox_styled(slide, x + Inches(0.8), y + Inches(0.4),
                           Inches(10), Inches(0.5),
                           desc, font_pt=12, color="666666")
    return None


def add_section_divider(slide, num, title, subtitle="",
                        dark_color="333333", accent_color="E89232"):
    """Transicao de secao com split panel."""
    SLIDE_W = Inches(13.333)
    SLIDE_H = Inches(7.5)

    # Left panel (60%)
    left_w = Inches(8)
    left = slide.shapes.add_shape(1, Emu(0), Emu(0), left_w, SLIDE_H)
    left.fill.solid()
    left.fill.fore_color.rgb = RGBColor.from_string(dark_color)
    left.line.fill.background()

    # Section number
    add_textbox_styled(slide, Inches(0.8), Inches(1.5), Inches(2), Inches(1),
                       f"{num:02d}", font_pt=72, color=accent_color, bold=True)

    # Title
    add_textbox_styled(slide, Inches(0.8), Inches(3.5), Inches(6.5), Inches(2),
                       title, font_pt=32, color="FFFFFF", bold=True)

    # Right panel (40%)
    right_x = left_w
    right_w = SLIDE_W - left_w
    right = slide.shapes.add_shape(1, right_x, Emu(0), right_w, SLIDE_H)
    right.fill.solid()
    right.fill.fore_color.rgb = RGBColor.from_string("F5F5F5")
    right.line.fill.background()

    # Subtitle on right
    if subtitle:
        add_textbox_styled(slide, right_x + Inches(0.5), Inches(1.5),
                           right_w - Inches(1), Inches(4.5),
                           subtitle, font_pt=14, color="555555")

    return None


def add_card_grid_2x2(slide, items, x=Inches(0.8), y=Inches(2),
                      w=Inches(11.5), h=Inches(4.5)):
    """4 items em grid 2x2."""
    gap = Inches(0.3)
    card_w = (w - gap) / 2
    card_h = (h - gap) / 2

    positions = [
        (x, y),
        (x + card_w + gap, y),
        (x, y + card_h + gap),
        (x + card_w + gap, y + card_h + gap),
    ]

    for i, item in enumerate(items[:4]):
        px, py = positions[i]
        add_card(slide, px, py, card_w, card_h,
                 title=item.get('title', ''),
                 body=item.get('body', ''))


def add_column_cards(slide, items, x=Inches(0.5), y=Inches(2),
                     total_w=Inches(12), h=Inches(4.5)):
    """N colunas paralelas."""
    n = len(items)
    if n == 0:
        return
    gap = Inches(0.2)
    card_w = (total_w - gap * (n - 1)) / n

    for i, item in enumerate(items):
        cx = x + i * (card_w + gap)
        add_card(slide, cx, y, card_w, h,
                 title=item.get('title', ''),
                 body=item.get('body', ''))


def add_footer_bar(slide, text, color="333333", font_color="AAAAAA"):
    """Rodape padrao do slide."""
    SLIDE_W = Inches(13.333)
    bar_h = Inches(0.4)
    bar_y = Inches(7.1)

    shape = slide.shapes.add_shape(1, Emu(0), bar_y, SLIDE_W, bar_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(color)
    shape.line.fill.background()

    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(8)
    run.font.name = "Arial"
    run.font.color.rgb = RGBColor.from_string(font_color)

    return shape


def force_font_on_all_runs(prs, font_name="Arial"):
    """Final pass: garantir fonte explicita em todos os runs."""
    ns = f'{{{NS_A}}}'
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for r in para._p.findall(f'{ns}r'):
                    rPr = r.find(f'{ns}rPr')
                    if rPr is None:
                        rPr = etree.SubElement(r, f'{ns}rPr', attrib={'lang': 'pt-BR'})
                    if rPr.find(f'{ns}latin') is None:
                        etree.SubElement(rPr, f'{ns}latin', attrib={'typeface': font_name})
                    if rPr.find(f'{ns}cs') is None:
                        etree.SubElement(rPr, f'{ns}cs', attrib={'typeface': font_name})


def set_slide_bg(slide, hex_color):
    """Setar background solido de um slide (dentro de cSld, nao sld)."""
    ns_p = f'{{{NS_P}}}'
    ns_a = f'{{{NS_A}}}'

    cSld = slide._element.find(f'{ns_p}cSld')
    old_bg = cSld.find(f'{ns_p}bg')
    if old_bg is not None:
        cSld.remove(old_bg)

    bg = etree.Element(f'{ns_p}bg')
    bgPr = etree.SubElement(bg, f'{ns_p}bgPr')
    solidFill = etree.SubElement(bgPr, f'{ns_a}solidFill')
    etree.SubElement(solidFill, f'{ns_a}srgbClr', attrib={'val': hex_color})
    etree.SubElement(bgPr, f'{ns_a}effectLst')

    cSld.insert(0, bg)
